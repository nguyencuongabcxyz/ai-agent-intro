"""
Build the "Anatomy of an AI Agent" presentation using the slide template.

Layouts used from template:
  0  - 1_Opening 1       (title slide with picture)
  2  - Cover 1            (section cover with picture + title)
  5  - Agenda             (agenda with numbered items)
  6  - Bullet text        (title + content)
  7  - 2x Bullet text     (title + 2 columns)
  8  - Subtitel 2x bullet (title + 2 subtitled columns)
  9  - Subtitel 3x bullet (title + 3 subtitled columns)
  10 - Titel slide        (title only - for big statements)
  51 - Questions          (Q&A slide)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import copy

TEMPLATE = "presentation/slide template.pptx"
OUTPUT = "presentation/Anatomy of an AI Agent.pptx"

prs = Presentation(TEMPLATE)

# Remove the 3 template example slides
for _ in range(len(prs.slides)):
    sldId = prs.slides._sldIdLst[0]
    rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    if rId:
        prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(sldId)

# --- Helper -----------------------------------------------------------

def get_layout(name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise ValueError(f"Layout '{name}' not found")


def add_bullet_slide(title, bullets, level=None):
    """Add a 'Bullet text' slide. bullets = list of strings.
       level = optional list of ints (0=top, 1=indented) same length as bullets."""
    slide = prs.slides.add_slide(get_layout("Bullet text"))
    slide.placeholders[0].text = title
    tf = slide.placeholders[16].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(16)
        lvl = 0 if level is None else level[i]
        p.level = lvl
    return slide


def add_two_col_slide(title, left_bullets, right_bullets):
    """Add a '2x Bullet text' slide."""
    slide = prs.slides.add_slide(get_layout("2x Bullet text"))
    slide.placeholders[0].text = title

    for ph_idx, bullets in [(18, left_bullets), (19, right_bullets)]:
        tf = slide.placeholders[ph_idx].text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.font.size = Pt(14)
    return slide


def add_sub2col_slide(title, sub_left, left_bullets, sub_right, right_bullets):
    """Subtitel 2x bullet text layout."""
    slide = prs.slides.add_slide(get_layout("Subtitel 2x bullet text"))
    slide.placeholders[0].text = title
    slide.placeholders[16].text = sub_left
    slide.placeholders[17].text = sub_right

    for ph_idx, bullets in [(20, left_bullets), (21, right_bullets)]:
        tf = slide.placeholders[ph_idx].text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.font.size = Pt(14)
    return slide


def add_sub3col_slide(title, sub1, col1, sub2, col2, sub3, col3):
    """Subtitel 3x bullet text layout."""
    slide = prs.slides.add_slide(get_layout("Subtitel 3x bullet text"))
    slide.placeholders[0].text = title
    slide.placeholders[16].text = sub1
    slide.placeholders[17].text = sub2
    slide.placeholders[19].text = sub3

    for ph_idx, bullets in [(20, col1), (21, col2), (22, col3)]:
        tf = slide.placeholders[ph_idx].text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.font.size = Pt(13)
    return slide


def add_title_slide(title):
    """Titel slide - big statement, title only."""
    slide = prs.slides.add_slide(get_layout("Titel slide"))
    slide.placeholders[0].text = title
    return slide


# ======================================================================
# SLIDES
# ======================================================================

# --- 1. Opening -------------------------------------------------------
slide = prs.slides.add_slide(get_layout("1_Opening 1"))
slide.placeholders[53].text = "Anatomy of an AI Agent"
slide.placeholders[13].text = "Building an AI agent from scratch to understand what's inside"

# --- 2. The Question ---------------------------------------------------
add_bullet_slide(
    "You use AI coding tools every day...",
    [
        "Cursor, Claude Code, Copilot, Devin — what's actually inside?",
        "They read your files, edit code, run tests, search the web",
        "Is this magic?",
        "",
        'By the end of this session, you\'ll know exactly how they work.',
        "Because we're going to build one. From scratch. In Python.",
    ],
)

# --- 3. Agenda ---------------------------------------------------------
add_bullet_slide(
    "Agenda",
    [
        "Stage 1  —  The Brain (LLM)",
        "Stage 2  —  The Agent Loop",
        "Stage 3  —  The Tools",
        "Stage 4  —  Extending Capabilities (MCP)",
        "Stage 5  —  Memory",
        "The Full Picture",
        "Should You Build an Agent?",
    ],
)

# ======================================================================
# STAGE 1: THE BRAIN
# ======================================================================

add_title_slide("Stage 1 — The Brain")

add_bullet_slide(
    "The LLM is the brain of the agent",
    [
        "An LLM (Large Language Model) is a text-in, text-out API",
        "You send it text, it sends back text. Nothing more.",
        "",
        "It CANNOT:",
        "Search the web",
        "Read or write files",
        "Run code",
        "Remember previous conversations",
        "",
        "A brain in a jar — it can think, but has no hands, no eyes, no memory.",
    ],
    level=[0, 0, 0, 0, 1, 1, 1, 1, 0, 0],
)

# Diagram slide for brain
add_sub2col_slide(
    "The Brain — How it works",
    "You ask:",
    [
        '"What files are in /src?"',
        "",
        "You send a message to the LLM API",
        "Just an HTTP request with JSON",
    ],
    "LLM responds:",
    [
        '"I don\'t know, I can\'t look at your files."',
        "",
        "It can only generate text based on training",
        "It has NO access to the outside world",
    ],
)

add_bullet_slide(
    "stage1_llm.py — Demo",
    [
        "The simplest possible AI application: a chatbot",
        "Sends a user message to the LLM, prints the response",
        "No tools. No loop. No memory.",
        "",
        "Live demo:",
        'Ask "What files are in this directory?"',
        "Watch it guess (or admit it can't)",
    ],
    level=[0, 0, 0, 0, 0, 1, 1],
)

add_bullet_slide(
    "Real-World Facts: The Brain",
    [
        "Claude Opus: ~$15/M input tokens, ~$75/M output tokens",
        "Every message to the brain has a price tag",
        "",
        "GPT-4o, Gemini, Claude — different brains, same interface",
        "Claude Code uses Claude. Cursor lets you pick. The brain is swappable.",
        "",
        "Claude Code's system prompt is ~12,000 tokens",
        "Before you even type, the brain has read a small novel of instructions",
    ],
)

add_title_slide("An LLM by itself is a chatbot.\nTo make it an agent, it needs a body.")

# ======================================================================
# STAGE 2: THE LOOP
# ======================================================================

add_title_slide("Stage 2 — The Agent Loop")

add_bullet_slide(
    "The agent loop turns a chatbot into an agent",
    [
        "Stage 1:  User  →  LLM  →  Response.  One shot. Done.",
        "Stage 2:  User  →  LLM  →  Tool  →  LLM  →  Tool  →  ...  →  Done.",
        "",
        "The LLM can now CALL TOOLS, and we LOOP until it's done.",
        "That loop is literally a while loop. That's the entire secret.",
    ],
)

# Agent loop diagram
add_bullet_slide(
    "The Agent Loop — Diagram",
    [
        "1.  Send messages to the LLM",
        "2.  LLM responds — does it want to call a tool?",
        "      YES → Execute the tool, add result to messages, go to step 1",
        "      NO  → We're done. Show the text response to the user.",
        "",
        "That's it. The entire 'agent' concept is this loop.",
        "Everything else is what you plug INTO this loop.",
    ],
    level=[0, 0, 1, 1, 0, 0, 0],
)

add_bullet_slide(
    "stage2_agent_loop.py — Demo",
    [
        "Added a while loop + 2 inline tools (web_search, read_file)",
        "The loop: ~10 lines of Python",
        "Stop condition: LLM returns text without tool calls",
        "",
        "Live demo:",
        "Ask something that requires a tool",
        "Watch the colored terminal output — see the loop iterate",
    ],
    level=[0, 0, 0, 0, 0, 1, 1],
)

add_bullet_slide(
    "Real-World Facts: The Loop",
    [
        "Claude Code — same loop: call Claude → check tool calls → execute → repeat",
        "Cursor — batches multiple tool calls in parallel",
        "Devin — adds a planning step before each loop iteration",
        "ChatGPT — chain-of-thought happens inside the brain, loop is the same",
        "",
        "A single prompt in Claude Code can trigger 50+ loop iterations",
        "One prompt → dozens of file reads, edits, grep searches, bash commands",
    ],
)

add_title_slide("Agent = LLM + Loop.\nEverything else is plugging things into this loop.")

# ======================================================================
# STAGE 3: THE TOOLS
# ======================================================================

add_title_slide("Stage 3 — The Tools")

add_bullet_slide(
    "Tools are the agent's hands",
    [
        "The LLM doesn't execute anything",
        "It ASKS us to execute, and we do it",
        "",
        "How a tool call works:",
        'LLM says: { "tool": "read_file", "args": { "path": "main.py" } }',
        'Our code: result = open("main.py").read()',
        "We feed the result back to the LLM",
        "The LLM continues with that new information",
    ],
    level=[0, 0, 0, 0, 1, 1, 1, 1],
)

add_sub3col_slide(
    "The 3-Part Tool Pattern",
    "Schema",
    [
        "JSON description of each tool",
        "Name, parameters, types",
        "Sent to LLM so it knows what's available",
        "Like a menu at a restaurant",
    ],
    "Registry",
    [
        "Lookup table",
        "Maps tool names to functions",
        '{"read_file": read_file_fn}',
        "Simple dictionary",
    ],
    "Dispatcher",
    [
        "Receives tool call from LLM",
        "Finds the function in registry",
        "Executes it with the arguments",
        "Returns result to the loop",
    ],
)

add_bullet_slide(
    "stage3_agent_tools.py — Demo",
    [
        "7 tools in a separate tools.py module",
        "read_file, write_file, list_files, web_search, fetch_webpage, run_bash, run_background",
        "",
        "The agent loop is UNCHANGED from Stage 2",
        "Same while loop, same stop condition",
        "Only the toolbox grew",
    ],
)

# Claude Code's tools
add_bullet_slide(
    "Real-World: Claude Code's Actual Tools",
    [
        "Read — Read files (like cat but smarter)",
        "Edit — Surgical string replacement in files",
        "Write — Create or overwrite files",
        "Bash — Run any shell command",
        "Grep — Search file contents (ripgrep)",
        "Glob — Find files by pattern",
        "Agent — Spawn a sub-agent (agent inside an agent!)",
        "WebSearch / WebFetch — Search and fetch from the web",
    ],
)

add_bullet_slide(
    "Interesting Facts: Tools",
    [
        "Claude Code has a permission system on tools",
        "Read runs freely. Bash and Write ask you first.",
        "The loop doesn't change — the tool harness adds a confirmation step",
        "",
        "Cursor's tools include apply_diff — generates a diff, tool applies it",
        "",
        "The LLM never sees your file system directly",
        "It only sees what tools RETURN. If read_file lies, the LLM believes it.",
        "Tools are the agent's only window to reality.",
    ],
)

add_title_slide("More tools = more capable agent.\nBut the loop is still the same.")

# ======================================================================
# STAGE 4: MCP
# ======================================================================

add_title_slide("Stage 4 — Extending Capabilities (MCP)")

add_bullet_slide(
    "What if tools didn't have to be hardcoded?",
    [
        "Stage 3: Tools are defined in tools.py — hardcoded at build time",
        "Stage 4: Tools are discovered from an MCP server at runtime",
        "",
        "MCP = Model Context Protocol",
        "Created by Anthropic, open-sourced late 2024",
        "Becoming the industry standard for AI tool connectivity",
    ],
)

add_sub2col_slide(
    "Before vs After MCP",
    "Before (Stage 3) — Hardcoded",
    [
        "Tools baked into the agent code",
        "read_file, write_file, web_search...",
        "Adding a tool = changing agent code",
        "Every agent rebuilds the same tools",
    ],
    "After (Stage 4) — Discovered",
    [
        "Agent asks: 'What tools do you have?'",
        "MCP server responds with tool list",
        "search_issues, create_ticket, get_board...",
        "Agent uses them — no code changes needed",
    ],
)

add_bullet_slide(
    "MCP = USB-C for AI",
    [
        "Plug in a server → get new tools",
        "Unplug it → tools disappear",
        "The laptop (agent) doesn't change",
        "",
        "The agent loop is IDENTICAL to Stage 2 and 3",
        "Only WHERE the tool calls go changed:",
        "Stage 3: tool call → local Python function",
        "Stage 4: tool call → external MCP server",
    ],
    level=[0, 0, 0, 0, 0, 0, 1, 1],
)

add_bullet_slide(
    "stage4_agent_mcp.py — Demo",
    [
        "Connects to an MCP server (e.g. Atlassian)",
        "Discovers tools dynamically at startup",
        "The agent can now search Jira, create tickets, etc.",
        "",
        "We didn't write those tools — the MCP server provides them",
        "The agent loop code? Still the same while loop.",
    ],
)

add_bullet_slide(
    "Real-World: MCP Ecosystem",
    [
        "Claude Code supports MCP natively — Jira, Slack, GitHub, databases",
        "Cursor supports MCP too",
        "1000+ community MCP servers (GitHub, Postgres, Notion, Figma...)",
        "",
        "An MCP server is just a process speaking a JSON protocol",
        "You could write one in ~50 lines of Python",
        "",
        "Claude Desktop uses MCP for local files, Google Drive, Slack",
    ],
)

add_title_slide("Don't build tools INTO the agent.\nBuild tools FOR agents.")

# ======================================================================
# STAGE 5: MEMORY
# ======================================================================

add_title_slide("Stage 5 — Memory")

add_bullet_slide(
    "Without memory, every turn is a first date",
    [
        "Stage 3 (no memory):",
        'Turn 1: [system, user] → LLM          "My name is Cuong"',
        'Turn 2: [system, user] → LLM          "What\'s my name?" → Forgets!',
        "",
        "Stage 5 (with memory):",
        "Turn 1: [system, user1] → LLM",
        "Turn 2: [system, user1, assistant1, user2] → LLM → Remembers!",
        "",
        '"Memory" is just keeping the message list across turns.',
    ],
    level=[0, 1, 1, 0, 0, 1, 1, 0, 0],
)

add_bullet_slide(
    "stage5_agent_session.py — Demo",
    [
        "The messages list PERSISTS across turns",
        "Every user message, LLM response, tool call, and result — kept",
        "",
        "Live demo:",
        'Stage 3: "My name is Cuong" → "What\'s my name?" → Forgets',
        'Stage 5: Same questions → Remembers!',
    ],
    level=[0, 0, 0, 0, 1, 1],
)

# Memory levels in real world
add_bullet_slide(
    "Real-World: Memory Goes Much Deeper",
    [
        "Conversation memory — keep message history (our Stage 5)",
        "Project context — load relevant files automatically (Cursor indexes codebase, Claude Code uses CLAUDE.md)",
        "Persistent memory — remember across sessions (Claude Code memory files, ChatGPT memory)",
        "RAG — search a knowledge base for relevant context (enterprise agents)",
    ],
)

add_bullet_slide(
    "Interesting Facts: Memory",
    [
        "Claude Code's context: 1M tokens (~750k words = entire Harry Potter series)",
        "It still compresses old messages when conversations get long",
        "",
        "Cursor indexes your codebase into embeddings for instant recall",
        "ChatGPT memory = a text file of facts, injected into system prompt",
        "",
        "The hard problem isn't storing memory — it's deciding what to FORGET",
        "Context windows have limits. Every token costs money.",
    ],
)

add_title_slide("Brain + Loop + Tools + Discovery + Memory\n= a full agent.")

# ======================================================================
# THE FULL PICTURE
# ======================================================================

add_title_slide("The Full Picture")

add_bullet_slide(
    "Agent Architecture — All the pieces",
    [
        "Memory (messages) ←→ Agent Loop (while...) ←→ Tools (local + MCP)",
        "                                    ↕",
        "                              LLM (Brain)",
        "",
        "We built this in ~200 lines of Python.",
        "Claude Code, Cursor, Devin — they're this, plus engineering.",
    ],
)

# Mapping to real products
add_sub3col_slide(
    "How real AI agents implement this",
    "Claude Code",
    [
        "Brain: Claude",
        "Loop: Yes",
        "Tools: Read, Edit, Bash, Grep...",
        "MCP: Supported",
        "Memory: CLAUDE.md + memory files",
    ],
    "Cursor",
    [
        "Brain: Configurable (Claude, GPT...)",
        "Loop: Yes + parallel tool calls",
        "Tools: apply_diff, terminal...",
        "MCP: Supported",
        "Memory: Codebase indexing",
    ],
    "ChatGPT / Devin",
    [
        "Brain: GPT-4o / Multiple",
        "Loop: Yes / Yes + planner",
        "Tools: code interpreter, browsing",
        "MCP: Not yet",
        "Memory: Conversation + facts",
    ],
)

# ======================================================================
# SHOULD YOU BUILD AN AGENT?
# ======================================================================

add_title_slide("Should You Build an Agent?")

add_title_slide("Probably not.")

add_bullet_slide(
    "Why not?",
    [
        "The agent loop is a SOLVED problem",
        "Claude Code, Cursor, OpenAI Agents SDK, LangGraph — they all do it well",
        "",
        "Building your own agent means rebuilding:",
        "The loop, permissions, error handling, memory",
        "Streaming, retry logic, context management, UI",
        "",
        "The loop is commodity infrastructure now.",
        "Don't rebuild the engine. Build what goes on top.",
    ],
    level=[0, 0, 0, 0, 1, 1, 0, 0, 0],
)

add_sub2col_slide(
    "The Trend: Build Skills, Not Agents",
    "Old thinking",
    [
        '"I\'ll build an agent from scratch for my use case"',
        "",
        "Months of work",
        "Maintain the loop, tools, memory yourself",
        "Re-invent what already exists",
    ],
    "New thinking",
    [
        '"I\'ll build SKILLS that plug into existing agents"',
        "",
        "MCP servers — your tools, your data",
        "CLAUDE.md / .cursorrules — your conventions",
        "Hooks — your workflows",
        "Custom prompts — your expertise",
    ],
)

add_bullet_slide(
    "What 'building skills' looks like",
    [
        "MCP servers — give agents access to YOUR systems (Jira, internal APIs, databases)",
        "CLAUDE.md / .cursorrules — teach agents YOUR team's conventions",
        "Hooks — automate workflows around agent actions",
        "Custom prompts / slash commands — package expertise into reusable instructions",
        "",
        "You don't build the car. You decide where to drive it.",
    ],
)

# ======================================================================
# CLOSING
# ======================================================================

add_title_slide("You now understand the engine.\nNext topic: Agent Skills.")

# Questions slide
slide = prs.slides.add_slide(get_layout("Questions"))

# --- Save ---
prs.save(OUTPUT)
print(f"Saved {OUTPUT} with {len(prs.slides)} slides")
